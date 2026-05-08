"""
Image extractor for multimodal RAG ingestion.

Extracts embedded images from:
  - PDF files   → via pdfplumber (renders each page, crops image bboxes)
  - PPTX files  → via python-pptx (reads picture shapes)
  - Standalone images (.png, .jpg, .jpeg, .webp, .gif, .bmp, .tiff)

Each extracted image is returned as a PIL.Image along with structured
metadata describing its origin (page number, slide number, image index,
dimensions).  Images smaller than a configurable threshold are skipped to
avoid wasting captioning budget on decorative bullets / logos.
"""

from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, List, Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Public data type
# ---------------------------------------------------------------------------

@dataclass
class ExtractedImage:
    """
    A single image taken from a document together with provenance metadata.

    The ``image`` field is a ``PIL.Image.Image`` object ready to be passed
    directly to the captioner.  All other fields are informational and will
    be stored verbatim in the chunk metadata.
    """

    image: object                   # PIL.Image.Image  (avoid hard import in type hint)
    page_or_slide: int              # 1-based page (PDF) or slide (PPTX) number; 0 for standalone
    image_index: int                # 0-based index within the page/slide
    source_path: str                # Absolute path of the source document
    width: int = 0
    height: int = 0
    extra_metadata: dict = field(default_factory=dict)

    def to_metadata_dict(self) -> dict:
        return {
            "chunk_type": "image_caption",
            "page": self.page_or_slide,
            "image_index": self.image_index,
            "source_doc": self.source_path,
            "image_width": self.width,
            "image_height": self.height,
            **self.extra_metadata,
        }


# ---------------------------------------------------------------------------
# Main extractor
# ---------------------------------------------------------------------------

class ImageExtractor:
    """
    Extracts images from PDF, PPTX, and standalone image files.

    Parameters
    ----------
    min_width / min_height
        Images whose pixel dimensions fall below these thresholds are
        silently skipped (decorative icons, bullets, etc.).
    max_images_per_page
        Hard cap on images extracted per page / slide to guard against
        pathological files with thousands of tiny images.
    """

    STANDALONE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".tif"}
    PDF_EXTENSIONS = {".pdf"}
    PPTX_EXTENSIONS = {".pptx", ".ppt"}

    def __init__(
        self,
        min_width: int = 80,
        min_height: int = 80,
        max_images_per_page: int = 20,
    ):
        self.min_width = min_width
        self.min_height = min_height
        self.max_images_per_page = max_images_per_page

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    def extract(self, file_path: str) -> List[ExtractedImage]:
        """
        Extract all qualifying images from *file_path*.

        Returns an empty list (and logs a warning) if the file type is
        unsupported or if the required optional library is not installed.
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext in self.STANDALONE_EXTENSIONS:
            return self._extract_standalone(file_path)
        elif ext in self.PDF_EXTENSIONS:
            return self._extract_from_pdf(file_path)
        elif ext in self.PPTX_EXTENSIONS:
            return self._extract_from_pptx(file_path)
        else:
            logger.debug(f"ImageExtractor: unsupported extension '{ext}' for {file_path}")
            return []

    # ------------------------------------------------------------------
    # Standalone image
    # ------------------------------------------------------------------

    def _extract_standalone(self, file_path: str) -> List[ExtractedImage]:
        try:
            from PIL import Image as PILImage

            img = PILImage.open(file_path).convert("RGB")
            w, h = img.size
            if w < self.min_width or h < self.min_height:
                logger.debug(f"Skipping standalone image (too small): {file_path}")
                return []

            logger.info(f"Standalone image loaded: {os.path.basename(file_path)} ({w}×{h})")
            return [
                ExtractedImage(
                    image=img,
                    page_or_slide=0,
                    image_index=0,
                    source_path=file_path,
                    width=w,
                    height=h,
                )
            ]
        except ImportError:
            logger.warning("Pillow is not installed — cannot load standalone images.")
            return []
        except Exception as e:
            logger.error(f"Failed to load image {file_path}: {e}")
            return []

    # ------------------------------------------------------------------
    # PDF extraction
    # ------------------------------------------------------------------

    def _extract_from_pdf(self, file_path: str) -> List[ExtractedImage]:
        """
        Use pdfplumber to iterate pages and extract embedded image bytes.

        pdfplumber exposes ``page.images`` — each entry is a dict with
        ``stream`` (raw bytes), ``x0``, ``y0``, ``x1``, ``y1``, etc.
        We decode the stream via PIL and apply dimension filtering.
        """
        try:
            import pdfplumber
            from PIL import Image as PILImage
        except ImportError:
            logger.warning(
                "pdfplumber or Pillow not installed — skipping PDF image extraction. "
                "Install with: pip install pdfplumber Pillow"
            )
            return []

        extracted: List[ExtractedImage] = []

        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    raw_images = page.images or []
                    img_count = 0

                    for img_meta in raw_images:
                        if img_count >= self.max_images_per_page:
                            logger.debug(
                                f"PDF page {page_num}: hit max_images_per_page limit"
                            )
                            break

                        try:
                            data = img_meta.get("stream")
                            if data is None:
                                continue

                            # ``stream`` is a pdfminer Stream object; get bytes
                            if hasattr(data, "get_data"):
                                raw_bytes = data.get_data()
                            elif isinstance(data, (bytes, bytearray)):
                                raw_bytes = bytes(data)
                            else:
                                raw_bytes = bytes(data)

                            pil_img = PILImage.open(io.BytesIO(raw_bytes)).convert("RGB")
                            w, h = pil_img.size

                            if w < self.min_width or h < self.min_height:
                                continue

                            extracted.append(
                                ExtractedImage(
                                    image=pil_img,
                                    page_or_slide=page_num,
                                    image_index=img_count,
                                    source_path=file_path,
                                    width=w,
                                    height=h,
                                    extra_metadata={
                                        "pdf_filter": img_meta.get("Filter", ""),
                                        "color_space": img_meta.get("ColorSpace", ""),
                                    },
                                )
                            )
                            img_count += 1

                        except Exception as img_err:
                            logger.debug(
                                f"PDF page {page_num}, image #{img_count}: "
                                f"could not decode — {img_err}"
                            )

                    if img_count:
                        logger.info(
                            f"PDF '{os.path.basename(file_path)}' page {page_num}: "
                            f"extracted {img_count} image(s)"
                        )

        except Exception as e:
            logger.error(f"pdfplumber failed on {file_path}: {e}", exc_info=True)

        logger.info(
            f"PDF extraction complete: {len(extracted)} image(s) from "
            f"'{os.path.basename(file_path)}'"
        )
        return extracted

    # ------------------------------------------------------------------
    # PPTX extraction
    # ------------------------------------------------------------------

    def _extract_from_pptx(self, file_path: str) -> List[ExtractedImage]:
        """
        Use python-pptx to iterate slides and pull picture shape blobs.

        Shapes with ``shape_type == MSO_SHAPE_TYPE.PICTURE`` have an
        ``image.blob`` attribute containing the raw image bytes.
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import PP_PLACEHOLDER
            from pptx.util import Emu
            from PIL import Image as PILImage
        except ImportError:
            logger.warning(
                "python-pptx or Pillow not installed — skipping PPTX image extraction. "
                "Install with: pip install python-pptx Pillow"
            )
            return []

        extracted: List[ExtractedImage] = []

        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, start=1):
                img_count = 0

                for shape in slide.shapes:
                    if img_count >= self.max_images_per_page:
                        break

                    # Check for picture shapes (MSO_SHAPE_TYPE.PICTURE == 13)
                    if not (hasattr(shape, "shape_type") and shape.shape_type == 13):
                        continue

                    try:
                        blob = shape.image.blob
                        pil_img = PILImage.open(io.BytesIO(blob)).convert("RGB")
                        w, h = pil_img.size

                        if w < self.min_width or h < self.min_height:
                            continue

                        extracted.append(
                            ExtractedImage(
                                image=pil_img,
                                page_or_slide=slide_num,
                                image_index=img_count,
                                source_path=file_path,
                                width=w,
                                height=h,
                                extra_metadata={
                                    "image_content_type": shape.image.content_type,
                                    "shape_name": shape.name,
                                },
                            )
                        )
                        img_count += 1

                    except Exception as shape_err:
                        logger.debug(
                            f"PPTX slide {slide_num}, shape '{shape.name}': "
                            f"could not extract image — {shape_err}"
                        )

                if img_count:
                    logger.info(
                        f"PPTX '{os.path.basename(file_path)}' slide {slide_num}: "
                        f"extracted {img_count} image(s)"
                    )

        except Exception as e:
            logger.error(f"python-pptx failed on {file_path}: {e}", exc_info=True)

        logger.info(
            f"PPTX extraction complete: {len(extracted)} image(s) from "
            f"'{os.path.basename(file_path)}'"
        )
        return extracted
